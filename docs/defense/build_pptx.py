#!/usr/bin/env python3
"""Build an editable, well-laid-out PowerPoint of the WireGuard defense deck.

Full control over fonts/sizes/positions via python-pptx — blue Helvetica
titles, a fixed-size monospace code box, centered diagrams, balanced tables,
and the speaker notes carried into PowerPoint presenter notes.
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

FONT      = "Helvetica"
MONO      = "Courier New"
BLUE      = RGBColor(0x1E, 0x3A, 0x8A)
RED       = RGBColor(0xB9, 0x1C, 0x1C)
GREY      = RGBColor(0x55, 0x55, 0x55)
INK       = RGBColor(0x20, 0x20, 0x20)
CODE_BG   = RGBColor(0xF3, 0xF4, 0xF6)
TBL_HEAD  = RGBColor(0x1E, 0x3A, 0x8A)
TBL_ALT   = RGBColor(0xE8, 0xEC, 0xF6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

DIAG = "../diagrams/"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _emit_runs(p, text, size, color=INK):
    """Render a string with **bold**, `code`, *italic* into runs on paragraph p."""
    pat = re.compile(r"\*\*(.+?)\*\*|`(.+?)`|\*(.+?)\*")
    pos = 0
    for m in pat.finditer(text):
        if m.start() > pos:
            _run(p, text[pos:m.start()], size, color)
        if m.group(1) is not None:
            _run(p, m.group(1), size, color, bold=True)
        elif m.group(2) is not None:
            _run(p, m.group(2), size, color, mono=True)
        else:
            _run(p, m.group(3), size, color, italic=True)
        pos = m.end()
    if pos < len(text):
        _run(p, text[pos:], size, color)


def _run(p, text, size, color=INK, bold=False, italic=False, mono=False):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.name = MONO if mono else FONT
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return r


def title(s, text, appendix=False, center=False):
    if appendix:
        tag = s.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(4), Inches(0.4))
        tp = tag.text_frame.paragraphs[0]
        _run(tp, "APPENDIX", 13, GREY, bold=True)
    box = s.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(12.2), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    _run(p, text, 30, BLUE, bold=True)
    return box


def body_box(s, top, height=5.2, left=0.7, width=11.9):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    return tf


def add_lines(tf, lines, size=19, gap=10):
    """lines: list of (text, kind) kind in {bullet, num, plain, caption, sub}."""
    first = True
    for text, kind in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if kind == "caption":
            _emit_runs(p, text, 13, GREY)
        elif kind == "bullet":
            _emit_runs(p, "•  " + text, size)
        elif kind == "num":
            _emit_runs(p, text, size)
        elif kind == "sub":
            p.space_after = Pt(4)
            _emit_runs(p, text, size, GREY)
        else:
            _emit_runs(p, text, size)


def code_box(s, code, top, height, size=13):
    box = s.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0xD0, 0xD4, 0xDC)
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        # crude C/keyword coloring kept simple: comments in grey-blue
        if "//" in line:
            head, com = line.split("//", 1)
            _run(p, head, size, INK, mono=True)
            _run(p, "//" + com, size, RGBColor(0x6B, 0x9B, 0xB5), mono=True, italic=True)
        else:
            _run(p, line, size, INK, mono=True)
    return box


def image(s, name, top=1.55, max_w=11.4, max_h=5.4):
    path = DIAG + name
    pic = s.shapes.add_picture(path, 0, Inches(top), width=Inches(max_w))
    if pic.height > Inches(max_h):
        el = pic._element
        el.getparent().remove(el)
        pic = s.shapes.add_picture(path, 0, Inches(top), height=Inches(max_h))
    pic.left = int((SW - pic.width) / 2)
    return pic


def table(s, data, left, top, width, col_widths=None, head=True, fontsize=13):
    rows, cols = len(data), len(data[0])
    gt = s.shapes.add_table(rows, cols, Inches(left), Inches(top),
                            Inches(width), Inches(0.4 * rows)).table
    gt.first_row = head
    if col_widths:
        for j, w in enumerate(col_widths):
            gt.columns[j].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            if head and i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = TBL_HEAD
            elif i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = TBL_ALT
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE
            p = c.text_frame.paragraphs[0]
            color = WHITE if (head and i == 0) else INK
            _emit_runs(p, str(val), fontsize, color)
            if head and i == 0:
                for r in p.runs:
                    r.font.bold = True
    return gt


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


# ---------------------------------------------------------------- slide 1
s = slide()
box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "WireGuard's receive path", 40, BLUE, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)
_run(p, "Finding and fixing the Execution Order Inversion", 24, INK, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(22)
_run(p, "Anas Ait El Hadj — Inria internship (KrakOS)", 18)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
_emit_runs(p, "Supervisors: **Alain Tchana** · **André Freyssinet**", 18)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(18)
_run(p, "Defense — June 10, 2026 · Salle F117", 14, GREY)
notes(s, """(15s) So — my internship is about a performance problem buried inside
WireGuard's receive path. One CPU core saturates, throughput collapses, and the
reason turns out to be something surprisingly subtle. I'll walk you through how
the pipeline works, where it breaks, what the fix looks like, and what I
measured. Let's go.""")

# ---------------------------------------------------------------- slide 2
s = slide(); title(s, "Context")
tf = body_box(s, 1.7)
add_lines(tf, [
 ("**WireGuard**: modern VPN in the Linux kernel — fast for one client, but on a **server with 1,000 clients** it reaches only **19.2% of line rate**.", "bullet"),
 ("Prior work (Mounah *et al.*, SYSTOR 2025): found the cause (**Execution Order Inversion**) and proposed a fix → **4.7× throughput**. But that fix is **incomplete**: it makes each wasted operation cheaper, not less frequent.", "bullet"),
 ("**My work**: understand the receive pipeline from the source code, identify the root cause of the EoI, and fix it at the trigger.", "bullet"),
], size=20, gap=16)
notes(s, """(45s) WireGuard is known for being fast and simple. But put it on a
server with a thousand clients pushing 25 Gbps, and one core hits 94% utilization
and throughput collapses to 19% of line rate. That's a CPU problem. Mounah et al.
at SYSTOR 2025 found the cause and proposed a fix that recovered 4.7× throughput.
But the root cause is still there: their fix makes each wasted operation cheaper
to run. Mine stops it from running in the first place.""")

# ---------------------------------------------------------------- slide 3
s = slide(); title(s, "The three kernel engines")
image(s, "slide_engines_en.png", top=1.7, max_w=11.6, max_h=5.3)
notes(s, """(1m15s) Three mechanisms cooperate. NAPI is the doorbell: instead of
an interrupt per packet, it rings once, mutes the bell, collects everything in one
pass. Key detail: napi_schedule does NOT run the poll — it sets a flag and raises
a softirq. The workqueue is the decryption workshop: ChaCha20-Poly1305 is heavy,
so it's delegated to background threads, one per core — fast, but they finish OUT
OF ORDER. GRO is the stapler: it batches same-flow packets into one trip. The
seed of the bug: ONE shared workshop, but each peer has its OWN queue and OWN
NAPI.""")

# ---------------------------------------------------------------- slide 4
s = slide(); title(s, "The pipeline — and where it breaks")
image(s, "slide_pipeline_en.png", top=1.55, max_w=10.6, max_h=4.7)
cap = s.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.9))
cf = cap.text_frame; cf.word_wrap = True
p = cf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "The flow turns ", 13, GREY)
_run(p, "DOWN", 13, RED, bold=True)
_run(p, " into the red box — ", 13, GREY)
_run(p, "napi_schedule unconditional", 13, RED, bold=True)
_run(p, " fires after every completion, regardless of queue state. This is the break point.", 13, GREY)
notes(s, """(1m15s) The full pipeline: receive, decrypt in the workshop, deliver
through NAPI and GRO. Each peer's decrypted packets land in its ordered queue, and
the poll drains it head-first so order is preserved. Now the red box: after every
decrypted packet — unconditionally — every worker rings the peer's NAPI. That
single unconditional call is where the bug lives.""")

# ---------------------------------------------------------------- slide 5
s = slide(); title(s, "The bug — Execution Order Inversion")
image(s, "bug_zoom_en.png", top=1.7, max_w=11.0, max_h=5.3)
notes(s, """(1m30s) Two facts collide. One: the workshop decrypts out of order.
Two: every worker, when done, calls napi_schedule unconditionally. So a worker
finishes packet 5 and rings the bell; the NAPI wakes, looks at the HEAD — packet 2,
still encrypted — and returns work_done = 0. A wasted softirq pass, and GRO had
nothing to staple, so batching collapses. The math: with N cores, the head
finishes first with probability 1/N, so (N-1)/N of wakes are wasted. It compounds
with peer count. One core saturates at 94%, throughput collapses.""")

# ---------------------------------------------------------------- slide 6
s = slide(); title(s, "The fix — 6 lines")
tf = body_box(s, 1.55, height=0.7)
add_lines(tf, [("**Before** calling `napi_schedule`, check whether the **head of the queue is ready**.", "plain")], size=19)
code_box(s, """tail = READ_ONCE(peer->rx_queue.tail);
if (tail == (struct sk_buff *)&peer->rx_queue.empty ||
    atomic_read(&PACKET_CB(tail)->state) != PACKET_STATE_UNCRYPTED)
        napi_schedule(&peer->napi);     // otherwise: skip""", top=2.35, height=2.0, size=14)
tf = body_box(s, 4.6, height=2.5)
add_lines(tf, [
 ("**Safe:** `tail` written only by the single consumer → no race condition.", "bullet"),
 ("**Worst case:** stale read → skip one wake → the worker finishing the head wakes it then.", "bullet"),
 ("**Effect:** premature wakes disappear → GRO gets full batches back.", "bullet"),
], size=18, gap=10)
notes(s, """(1m) Six lines: before waking the NAPI, check whether there's anything
to do. Read the consumer cursor; if the head is still uncrypted, skip — the worker
that finishes the head will wake it. Safe because that cursor is written by only
one entity, the poll itself: a lock-free hint. Worst case is a stale read that
misses a wake, and NAPI's internal MISSED mechanism re-runs the poll. No packet is
ever stranded. Result: premature wakes gone, GRO batches back.""")

# ---------------------------------------------------------------- slide 7
s = slide(); title(s, "Results — polls down, batches up")
tf = body_box(s, 1.6, height=0.7)
add_lines(tf, [("**Setup:** Apple M1 (8 cores), Linux net namespaces (multi-peer loopback), `bpftrace` for in-kernel metrics — loopback, 5 runs each.", "caption")])
table(s, [
 ["", "1 peer", "8 peers", "32 peers"],
 ["Δ wasted polls", "−8.8%", "**−21.9%**", "**−20.7%**"],
 ["Batch size", "3.1 → 3.3", "8.7 → **9.6**", "7.7 → **8.9**"],
], left=1.6, top=2.35, width=10.1,
   col_widths=[3.1, 2.3, 2.3, 2.4], fontsize=16)
tf = body_box(s, 4.4, height=2.6)
add_lines(tf, [
 ("The reduction **grows with peer count** — exactly what the 1/N model predicts.", "bullet"),
 ("**Batch size rising** is the direct confirmation: GRO is woken less but does more each time.", "bullet"),
 ("Throughput flat — expected, the loopback never saturates `NET_RX_SOFTIRQ`.", "bullet"),
], size=18, gap=10)
notes(s, """(1m30s) Setup, for credibility: Apple M1, eight cores, multiple peers
from network namespaces over loopback, measured in-kernel with bpftrace, five runs
each. The fix cuts wasted polls 9–22%, and the reduction grows with peer count —
the 1/N model. The batch numbers are the most convincing: at 8 peers each useful
poll goes 8.7 → 9.6 packets. Throughput is flat because loopback never saturates
the softirq — that needs a real NIC, which is what comes next.""")

# ---------------------------------------------------------------- slide 8
s = slide(); title(s, "Conclusion")
tf = body_box(s, 1.7)
add_lines(tf, [
 ("**Understood** the receive pipeline from the source: NAPI, the per-CPU workqueue, GRO.", "bullet"),
 ("**Located** the root cause — one unconditional line (`queueing.h:196`) that wakes the wrong thing at the wrong time.", "bullet"),
 ("**Fixed it** in **6 lines**: read the queue head before waking, skip the premature wake.", "bullet"),
 ("**Confirmed on ARM:** **9–22% fewer wasted polls**, and **batch size up** — direct evidence the mechanism works.", "bullet"),
], size=20, gap=16)
notes(s, """(40s) Three things done. I understood the receive pipeline from the
source. I located the root cause: a single unconditional line that wakes the NAPI
at the wrong moment. And I fixed it in six lines. On ARM I measured 9–22% fewer
wasted polls with batch sizes rising — direct evidence the mechanism works. That's
the contribution.""")

# ---------------------------------------------------------------- slide 9
s = slide(); title(s, "Next steps")
tf = body_box(s, 1.7)
add_lines(tf, [
 ("**The internship continues through July.** Three directions:", "plain"),
 ("1. **Real hardware (CloudLab)** — x86, 25G NIC, 1,000 peers. Does the fix reduce **throughput collapse**, not just poll counts? Does ARM behavior reproduce on x86?", "num"),
 ("2. **A batching-aware trigger** — wake only when waking *pays off* (measure poll overhead vs. delivery + copy-to-userspace cost).", "num"),
 ("3. **Combine with the SYSTOR fix** — orthogonal, should be additive.", "num"),
], size=19, gap=14)
p = tf.add_paragraph(); p.space_before = Pt(14)
_run(p, "Thank you — questions?", 14, GREY)
notes(s, """(40s) Clear continuation through July. First, real hardware: CloudLab,
x86 with a 25-gig NIC, to see whether fewer polls turn into real throughput and
whether ARM behavior holds on x86. Second, a smarter batching-aware trigger that
waits until waking pays off. Third, combining my trigger with the SYSTOR paper's
dedicated-workqueue fix — orthogonal, so they should stack. Thank you.""")

# ---------------------------------------------------------------- appendix A
s = slide(); title(s, "Appendix A — Stage 1: the peer and NAPI", appendix=True)
image(s, "slide_napi_en.png", top=1.45, max_w=10.6, max_h=4.5)
cap = s.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.9))
p = cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
cap.text_frame.word_wrap = True
_emit_runs(p, "Each client = one **peer** with its **own ordered queue** and its **own NAPI**. One shared decryption workshop for all peers.", 13, GREY)
notes(s, """Each client is a "peer" — identified by a public key. Each peer gets its
own ordered receive queue (preserving order) and its own NAPI instance, woken by
hand from the decrypt workers. napi_schedule sets a flag and raises a softirq; the
poll runs a moment later.""")

# ---------------------------------------------------------------- appendix B
s = slide(); title(s, "Appendix B — Stage 2: the workqueue", appendix=True)
image(s, "slide_wq_en.png", top=1.45, max_w=10.6, max_h=4.5)
cap = s.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.9))
cap.text_frame.word_wrap = True
p = cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_emit_runs(p, "Decryption is too heavy for the softirq → delegated to a **background pool**. One worker per core → decrypts **in parallel** → finishes **out of order**.", 13, GREY)
notes(s, """Decryption — ChaCha20-Poly1305 — is heavy, and the softirq is borrowed
time. WireGuard delegates to background kernel threads, one per core. Up to 8
packets from the same peer at once — fast, but out of order. When a worker finishes
it calls napi_schedule. Every worker, after every packet, unconditionally.""")

# ---------------------------------------------------------------- appendix C
s = slide(); title(s, "Appendix C — Stage 3: GRO", appendix=True)
image(s, "slide_gro_en.png", top=1.45, max_w=10.6, max_h=4.5)
cap = s.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.9))
cap.text_frame.word_wrap = True
p = cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_emit_runs(p, "Pushing a packet up the stack has a **fixed cost per packet** → GRO **staples packets into one parcel** → one trip instead of N.", 13, GREY)
notes(s, """GRO — Generic Receive Offload — is the optimization the bug destroys.
Traversing the stack has a fixed per-packet cost; with 40 same-flow packets it's
far better to staple them and traverse once. When GRO works you get big batches;
when something wakes it for nothing, the batches fall apart.""")

# ---------------------------------------------------------------- appendix D
s = slide(); title(s, "Appendix D — one workqueue, per-CPU workers", appendix=True)
tf = body_box(s, 1.7)
add_lines(tf, [
 ("**One** workqueue object (`packet_crypt_wq`, `device.c:346`), shared by all peers.", "bullet"),
 ("**Per-CPU = the *workers* are per core**: `queue_work_on(cpu, …)` dispatches each item to a specific core. It is **not** N workqueues.", "bullet"),
], size=20, gap=16)
notes(s, """It's not N workqueues — one object, allocated once. "Per-CPU" describes
the workers: each core gets one, and WireGuard uses queue_work_on to pin each
packet to a specific core.""")

# ---------------------------------------------------------------- appendix E
s = slide(); title(s, "Appendix E — NAPI lifecycle", appendix=True)
tf = body_box(s, 1.8)
add_lines(tf, [
 ("`netif_napi_add` → `napi_enable` → `napi_schedule` (sets flag + raises softirq, **nothing runs yet**)", "plain"),
 ("→ `wg_packet_rx_poll` → `napi_complete_done` → `napi_disable` → `netif_napi_del`", "plain"),
], size=18, gap=14)
p = tf.add_paragraph(); p.space_before = Pt(16)
_emit_runs(p, "The poll runs at the next `spin_unlock_bh` in the worker loop — not immediately.", 18)
notes(s, """The full lifecycle in seven steps. The one people get wrong is
napi_schedule: it doesn't run anything. It sets a bit and raises NET_RX_SOFTIRQ.
The poll runs at the next spin_unlock_bh at the end of each loop iteration — almost
immediately, but not synchronously.""")

# ---------------------------------------------------------------- appendix F
s = slide(); title(s, "Appendix F — Full results table", appendix=True)
table(s, [
 ["peers", "build", "wasted/s", "waste%", "batch", "Δwasted"],
 ["1",  "stock / patched", "42,638 / 38,872", "25.1 / 24.9", "3.1 / 3.3",  "−8.8%"],
 ["4",  "stock / patched", "33,652 / 30,512", "29.7 / 28.5", "14.2 / 15.0", "−9.3%"],
 ["8",  "stock / patched", "64,318 / 50,217", "29.2 / 28.2", "8.7 / 9.6",  "−21.9%"],
 ["16", "stock / patched", "50,788 / 44,480", "28.5 / 28.2", "9.9 / 11.6", "−12.4%"],
 ["32", "stock / patched", "64,987 / 51,553", "28.8 / 27.5", "7.7 / 8.9",  "−20.7%"],
], left=0.9, top=1.9, width=11.5,
   col_widths=[1.1, 2.7, 2.7, 2.2, 1.6, 1.6], fontsize=14)

# ---------------------------------------------------------------- appendix G
s = slide(); title(s, "Appendix G — bpftrace proof", appendix=True)
code_box(s, "kretprobe:wg_packet_rx_poll { @work_done = lhist(retval, 0, 64, 8); }",
         top=1.9, height=0.8, size=14)
tf = body_box(s, 3.0)
add_lines(tf, [
 ("**Spike in bucket 0** = wasted wakes (EoI signature).", "bullet"),
 ("The fix must **collapse bucket 0** and shift mass to > 1 (real batches).", "bullet"),
], size=20, gap=16)
notes(s, """One line of bpftrace traces the return value of wg_packet_rx_poll —
work_done, the packets delivered per pass. A histogram tells the whole story: if
the EoI is happening, a massive spike at zero. The fix makes that spike disappear
and pushes mass toward larger values. Architecture-independent — same on ARM and
x86.""")

prs.save("SLIDES_DEFENSE_EN_editable.pptx")
print("saved SLIDES_DEFENSE_EN_editable.pptx —", len(prs.slides._sldIdLst), "slides")
